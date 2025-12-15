from astrbot.api.event import filter, AstrMessageEvent          # pyright: ignore[reportMissingImports]
from astrbot.api.star import Context, Star, register            # pyright: ignore[reportMissingImports]
from astrbot.api.provider import ProviderRequest                # pyright: ignore[reportMissingImports]
from astrbot.core.provider.provider import RerankProvider        # pyright: ignore[reportMissingImports]
from astrbot.api.platform import MessageType                    # pyright: ignore[reportMissingImports]
import astrbot.api.message_components as Comp                   # pyright: ignore[reportMissingImports] 
from astrbot.api.all import *                                   # pyright: ignore[reportMissingImports] 
from astrbot.api import logger                                  # pyright: ignore[reportMissingImports]

# 导入配置相关模块
from pathlib import Path
import time
import os
import asyncio
import sqlite3

# 导入知识库相关模块
from astrbot.core.knowledge_base.chunking.recursive import RecursiveCharacterChunker
from astrbot.core.db.vec_db.faiss_impl.vec_db import FaissVecDB

# 导入文件处理相关模块
from pdfminer.high_level import extract_text
import docx2txt
import pandas as pd
from docx import Document
from pptx import Presentation
from typing import Dict, Optional
import chardet

# 使用字典存储支持的文件类型和对应的处理函数
SUPPORTED_EXTENSIONS: Dict[str, str] = {
    # 文档格式
    "pdf": "read_pdf_to_text",
    "docx": "read_docx_to_text",
    "doc": "read_docx_to_text",
    "rtf": "read_txt_to_text",
    "odt": "read_txt_to_text",

    # 电子表格
    "xlsx": "read_excel_to_text",
    "xls": "read_excel_to_text",
    "ods": "read_excel_to_text",
    "csv": "read_csv_to_text",

    # 演示文稿
    "pptx": "read_pptx_to_text",
    "ppt": "read_pptx_to_text",
    "odp": "read_pptx_to_text",

    # 编程语言源代码
    "py": "read_txt_to_text",
    "java": "read_txt_to_text",
    "cpp": "read_txt_to_text",
    "c": "read_txt_to_text",
    "h": "read_txt_to_text",
    "hpp": "read_txt_to_text",
    "cs": "read_txt_to_text",
    "js": "read_txt_to_text",
    "ts": "read_txt_to_text",
    "php": "read_txt_to_text",
    "rb": "read_txt_to_text",
    "go": "read_txt_to_text",
    "rs": "read_txt_to_text",
    "swift": "read_txt_to_text",
    "kt": "read_txt_to_text",
    "scala": "read_txt_to_text",
    "sh": "read_txt_to_text",
    "bash": "read_txt_to_text",
    "ps1": "read_txt_to_text",
    "bat": "read_txt_to_text",
    "cmd": "read_txt_to_text",
    "vbs": "read_txt_to_text",

    # 标记语言
    "html": "read_txt_to_text",
    "htm": "read_txt_to_text",
    "xml": "read_txt_to_text",
    "json": "read_txt_to_text",
    "yaml": "read_txt_to_text",
    "yml": "read_txt_to_text",
    "md": "read_txt_to_text",
    "markdown": "read_txt_to_text",

    # 配置文件
    "ini": "read_txt_to_text",
    "cfg": "read_txt_to_text",
    "conf": "read_txt_to_text",
    "properties": "read_txt_to_text",
    "env": "read_txt_to_text",

    # 数据库/查询
    "sql": "read_txt_to_text",

    # 其他文本格式
    "txt": "read_txt_to_text",
    "log": "read_txt_to_text",
    "": "read_txt_to_text",  # 无扩展名文件

    # 构建/项目文件
    "toml": "read_txt_to_text",
    "lock": "read_txt_to_text",
    "gitignore": "read_txt_to_text",

    # 网络相关
    "url": "read_txt_to_text",
    "webloc": "read_txt_to_text",
}

def get_file_type(file_path: str) -> Optional[str]:
    """安全获取文件扩展名（优先MIME检测，后备扩展名）"""
    
    # 首先检查文件是否存在
    if not os.path.isfile(file_path):
        raise FileNotFoundError
    
    try:
        # 方案1：使用python-magic（推荐）
        import magic
        mime = magic.from_file(file_path, mime=True)
        
        # 处理常见的MIME类型映射
        mime_to_ext = {
            'application/pdf': 'pdf',
            'image/jpeg': 'jpg',
            'image/png': 'png',
            'image/gif': 'gif',
            'text/plain': 'txt',
            'application/zip': 'zip',
            'application/x-rar-compressed': 'rar',
            'application/x-tar': 'tar',
            'application/gzip': 'gz'
        }
        
        # 处理常见的MIME类型映射
        if 'vnd.openxmlformats-officedocument' in mime:
            # 提取具体的Office文档类型
            if 'wordprocessingml' in mime:
                return 'docx'
            elif 'spreadsheetml' in mime:
                return 'xlsx'
            elif 'presentationml' in mime:
                return 'pptx'
        
        # 检查映射表
        if mime in mime_to_ext:
            return mime_to_ext[mime]
        
        # 通用处理：从MIME类型中提取扩展名
        if '/' in mime:
            mime_type = mime.split("/")[-1]
            # 处理复合类型如vnd.ms-excel
            if mime_type.startswith('vnd.'):
                mime_type = mime_type[4:]
            if mime_type.startswith('x-'):
                mime_type = mime_type[2:]
            return mime_type
        
        return mime
    
    except ImportError:
        # 方案2：后备使用扩展名
        ext = os.path.splitext(file_path)[1]
        if ext:
            return ext[1:].lower()  # 去掉点号并转为小写
        else:
            raise ImportError


def complete_filename(file_path: str) -> str:
    """补全文件名（如果缺少扩展名则自动添加）"""
    if not os.path.isfile(file_path):
        return file_path
    
    # 如果已经有扩展名，直接返回
    if os.path.splitext(file_path)[1]:
        return file_path
    
    # 获取文件类型并补全扩展名
    file_type = get_file_type(file_path)
    if file_type:
        return f"{file_path}.{file_type}"
    
    return file_path  # 无法确定类型，返回原文件名


def read_csv_to_text(file_path: str) -> str:
    """读取CSV文件并返回格式化的文本"""
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        raise RuntimeError(f"读取CSV文件失败: {str(e)}")


def read_pdf_to_text(file_path: str) -> str:
    """使用pdfminer.six提取PDF文本（效果更好）"""
    try:
        return extract_text(file_path)
    except Exception as e:
        raise RuntimeError(f"读取PDF文件失败: {str(e)}")


def convert_doc_to_docx(doc_file: str, docx_file: str) -> None:
    """将doc文档转为docx文档"""
    try:
        doc = Document(doc_file)
        doc.save(docx_file)
    except Exception as e:
        raise RuntimeError(f"转换DOC到DOCX失败: {str(e)}")


def read_docx_to_text(file_path: str) -> str:
    """读取DOCX或DOC文件内容并返回文本"""
    try:
        # 统一处理路径
        file_path = os.path.abspath(file_path)

        if file_path.lower().endswith(".doc"):
            # 处理DOC文件
            file_dir, file_name = os.path.split(file_path)
            file_base = os.path.splitext(file_name)[0]
            docx_file = os.path.join(file_dir, f"{file_base}.docx")

            # 转换DOC到DOCX
            convert_doc_to_docx(file_path, docx_file)

            # 处理转换后的文件
            text = docx2txt.process(docx_file)

            # 删除临时转换的文件
            try:
                os.remove(docx_file)
            except:
                pass
        else:
            # 直接处理DOCX文件
            text = docx2txt.process(file_path)

        return text
    except Exception as e:
        raise RuntimeError(f"读取Word文件失败: {str(e)}")


def read_excel_to_text(file_path: str) -> str:
    """读取Excel文件内容并返回文本"""
    try:
        excel_file = pd.ExcelFile(file_path)
        text_list = []

        for sheet_name in excel_file.sheet_names:
            df = excel_file.parse(sheet_name)
            text = df.to_string(index=False)
            text_list.append(f"=== {sheet_name} ===\n{text}")

        return "\n\n".join(text_list)
    except Exception as e:
        raise RuntimeError(f"读取Excel文件失败: {str(e)}")


def read_pptx_to_text(file_path: str) -> str:
    """读取PPTX文件内容并返回文本"""
    try:
        prs = Presentation(file_path)
        text_list = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    text_frame = shape.text_frame
                    if text_frame.text.strip():
                        slide_text.append(text_frame.text.strip())

            if slide_text:  # 只添加有内容的幻灯片
                text_list.append("\n".join(slide_text))

        return "\n\n".join(text_list)
    except Exception as e:
        raise RuntimeError(f"读取PPTX文件失败: {str(e)}")


def read_txt_to_text(file_path: str) -> str:
    """读取文本文件，自动检测编码"""
    try:
        with open(file_path, "rb") as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
        return raw_data.decode(encoding)
    except Exception as e:
        raise RuntimeError(f"读取文本文件失败: {str(e)}")


def read_any_file_to_text(file_path: str) -> str:
    """
    根据文件扩展名自动选择适当的读取函数
    返回文件内容文本或错误信息
    """
    try:
        # 修复路径编码问题
        if isinstance(file_path, bytes):
            try:
                file_path = file_path.decode('utf-8')
            except UnicodeDecodeError:
                file_path = file_path.decode('latin1')
        
        # 标准化路径（处理反斜杠和特殊字符）
        file_path = os.path.abspath(os.path.normpath(file_path))
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return f"文件不存在: {file_path}"
            
        # 获取文件扩展名（小写，不带点）
        file_ext = get_file_type(file_path)
        if file_ext is None or file_ext == "":
            file_ext = os.path.splitext(file_path)[1][1:].lower()
        if not file_ext:
            file_ext = "txt"  # 默认文本类型

        # 后续处理逻辑保持不变...
        func_name = SUPPORTED_EXTENSIONS.get(file_ext)
        if not func_name:
            return f"不支持 {file_ext} 格式"
            
        # 使用函数映射
        func_map = {
            "read_pdf_to_text": read_pdf_to_text,
            "read_docx_to_text": read_docx_to_text,
            "read_excel_to_text": read_excel_to_text,
            "read_pptx_to_text": read_pptx_to_text,
            "read_txt_to_text": read_txt_to_text,
            "read_csv_to_text": read_csv_to_text,
        }
        
        func = func_map.get(func_name)
        if func is None:
            return f"找不到处理 {file_ext} 文件的函数"
            
        return func(file_path)
        
    except Exception as e:
        return f"读取文件时出错: {str(e)}"


@register("astrbot_plugin_file_reader_pro", "zz6zz666", "一个将文件内容高效传给llm的插件（增强版）", "3.0.0")
class AstrbotPluginFileReaderPro(Star):
    PLUGIN_ID = "astrbot_plugin_file_reader_pro"
    
    def __init__(self, context: Context, config):
        super().__init__(context)
        self.file_name = ""
        self.file_dir = ""
        self.content = ""
        self.embedding_provider = None
        self.rerank_provider = None
        self.file_upload_time = None  # 文件上传时间
        self.config = self._load_config(config)  # 加载配置
        
        # 初始化所有配置项为类属性
        self.chunk_size = self.config.get("chunk_size", 512)
        self.chunk_overlap = self.config.get("chunk_overlap", 100)
        self.retrieve_top_k = self.config.get("retrieve_top_k", 5)
        self.fetch_k = self.config.get("fetch_k", 20)
        self.enable_rerank = self.config.get("enable_rerank", True)
        self.file_retention_time = self.config.get("file_retention_time", 60)  # 60分钟
        self.max_file_size = self.config.get("max_file_size", 100)  # 100MB
        self.file_max_rounds = self.config.get("file_max_rounds", 5)  # 文件最大使用轮数
        self.supported_file_types = self.config.get("supported_file_types", list(SUPPORTED_EXTENSIONS.keys()))
        self.rerank_provider_id = self.config.get("rerank_provider_id", "")  # 重排序模型服务商
        self.embedding_provider_id = self.config.get("embedding_provider_id", "")  # Embedding服务提供商
        self.cleanup_interval = self.config.get("cleanup_interval", 15)  # 清理间隔（分钟）
        self.enable_group_file_processing = self.config.get("enable_group_file_processing", True)  # 是否启用群文件处理
        self.enabled_groups = self.config.get("enabled_groups", [])  # 启用的群列表
        self.injection_type = self.config.get("injection_type", "system")  # 文件内容注入类型
        self.system_context_keep_rounds = self.config.get("system_context_keep_rounds", 2) # 系统上下文保留轮数
        
        # 初始化数据目录
        self._base_dir = Path(__file__).resolve().parent
        self._data_dir = self._resolve_data_dir()
        
        # 使用配置初始化分块器
        self.chunker = RecursiveCharacterChunker(chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap)
        
        # 当前活跃的会话和对话信息
        self.current_session_id = None
        self.current_conversation_id = None
        
        # 向量数据库实例字典，键为(session_id, conversation_id, file_name)
        self.vec_dbs = {}
        
        # 定期清理任务相关
        self._cleanup_task = None
        self._cleanup_interval = None
        
        # 初始化文件使用次数数据库连接
        self._init_file_rounds_db()
    
    def _load_config(self, config=None):
        """加载插件配置"""
        # 使用框架传入的配置
        if not config:
            # 如果配置不存在，使用默认值
            config = {}
            logger.info("未找到配置文件，使用默认配置")
        else:
            logger.info("配置加载成功")
        
        # 合并默认配置
        default_config = {
            "chunk_size": 512,
            "chunk_overlap": 100,
            "retrieve_top_k": 5,
            "fetch_k": 20,
            "enable_rerank": True,
            "file_retention_time": 60,  # 60分钟
            "max_file_size": 100,  # 100MB
            "file_max_rounds": 5,  # 文件最大使用轮数
            "supported_file_types": list(SUPPORTED_EXTENSIONS.keys()),
            "rerank_provider_id": "",  # 重排序模型服务商
            "embedding_provider_id": ""  # Embedding服务提供商
        }
        
        # 用默认配置填充缺失的配置项
        for key, value in default_config.items():
            if key not in config:
                config[key] = value
        
        return config
    
    def _resolve_data_dir(self) -> Path:
        """优先使用 StarTools 数据目录，失败时退回到 AstrBot/data/plugin_data 下。"""
        fallback_dir = self._base_dir.parent.parent / "plugin_data" / self.PLUGIN_ID
        try:
            from astrbot.api.star import StarTools
            preferred_raw = StarTools.get_data_dir(self.PLUGIN_ID)
            if preferred_raw:
                preferred_path = Path(preferred_raw)
                preferred_path.mkdir(parents=True, exist_ok=True)
                return preferred_path
        except Exception as exc:
            logger.warning(f"[文件读取插件] 创建数据目录失败({exc})，退回 fallback：{fallback_dir}")
        
        fallback_dir.mkdir(parents=True, exist_ok=True)
        return fallback_dir
        
    def _init_file_rounds_db(self):
        """初始化文件使用次数数据库"""
        # 数据库文件路径
        self._db_path = self._data_dir / "file_rounds.db"
        
        try:
            # 保持数据库连接打开以提高性能
            self._db_conn = sqlite3.connect(self._db_path)
            cursor = self._db_conn.cursor()
            
            # 创建文件使用次数表
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS file_rounds (
                    session_id TEXT NOT NULL,
                    conversation_id TEXT NOT NULL,
                    file_name TEXT NOT NULL,
                    rounds INTEGER DEFAULT 0,
                    PRIMARY KEY (session_id, conversation_id, file_name)
                )
            ''')
            
            # 创建索引以提高查询性能
            cursor.execute('''
                CREATE INDEX IF NOT EXISTS idx_file_rounds_session_conversation ON file_rounds (session_id, conversation_id)
            ''')
            
            self._db_conn.commit()
            logger.info(f"文件使用次数数据库初始化成功，路径：{self._db_path}")
        except Exception as e:
            logger.error(f"初始化文件使用次数数据库失败: {str(e)}")
            self._db_conn = None
            
    async def _start_periodic_cleanup(self):
        """启动定期清理任务"""
        # 使用类属性获取清理间隔
        cleanup_interval_seconds = self.cleanup_interval * 60  # 转换为秒
        
        async def cleanup_loop():
            while True:
                await asyncio.sleep(cleanup_interval_seconds)
                await self._cleanup_expired_files()
        
        self._cleanup_task = asyncio.create_task(cleanup_loop())
        logger.info(f"已启动定期清理任务，间隔：{self.cleanup_interval}分钟")
        
    async def _stop_periodic_cleanup(self):
        """停止定期清理任务"""
        if self._cleanup_task:
            self._cleanup_task.cancel()
            try:
                await self._cleanup_task
            except asyncio.CancelledError:
                pass
            logger.info("已停止定期清理任务")
            
    async def _cleanup_expired_files(self):
        """清理所有过期文件"""
        logger.info("开始执行定期清理任务")
        
        # 遍历所有向量数据库实例
        keys_to_cleanup = []
        for (session_id, conversation_id, file_name), vec_db in list(self.vec_dbs.items()):
            if self._is_file_expired(session_id, conversation_id, file_name):
                keys_to_cleanup.append((session_id, conversation_id, file_name))
        
        # 清理过期文件
        if keys_to_cleanup:
            logger.info(f"发现 {len(keys_to_cleanup)} 个过期文件，开始清理")
            for session_id, conversation_id, file_name in keys_to_cleanup:
                await self.cleanup(session_id, conversation_id, file_name)
        else:
            logger.info("未发现过期文件")
            
    def _get_file_rounds(self, session_id: str, conversation_id: str, file_name: str) -> int:
        """获取文件的使用次数"""
        if not self._db_conn:
            return 0
            
        try:
            cursor = self._db_conn.cursor()
            
            cursor.execute(
                "SELECT rounds FROM file_rounds WHERE session_id=? AND conversation_id=? AND file_name=?",
                (session_id, conversation_id, file_name)
            )
            
            result = cursor.fetchone()
            return result[0] if result else 0
        except Exception as e:
            logger.error(f"获取文件使用次数失败: {str(e)}")
            return 0
            
    def _increment_file_rounds(self, session_id: str, conversation_id: str, file_name: str):
        """增加文件的使用次数"""
        if not self._db_conn:
            return
            
        try:
            cursor = self._db_conn.cursor()
            
            # 使用upsert操作，不存在则插入，存在则更新
            cursor.execute(
                "INSERT OR REPLACE INTO file_rounds (session_id, conversation_id, file_name, rounds) VALUES (?, ?, ?, COALESCE((SELECT rounds + 1 FROM file_rounds WHERE session_id=? AND conversation_id=? AND file_name=?), 1))",
                (session_id, conversation_id, file_name, session_id, conversation_id, file_name)
            )
            
            self._db_conn.commit()
        except Exception as e:
            logger.error(f"增加文件使用次数失败: {str(e)}")
            
    def _delete_file_rounds(self, session_id: str, conversation_id: str, file_name: str = None):
        """删除文件的使用次数记录"""
        if not self._db_conn:
            return
            
        try:
            cursor = self._db_conn.cursor()
            
            if file_name:
                # 删除特定文件的记录
                cursor.execute(
                    "DELETE FROM file_rounds WHERE session_id=? AND conversation_id=? AND file_name=?",
                    (session_id, conversation_id, file_name)
                )
            elif conversation_id:
                # 删除整个对话的所有文件记录
                cursor.execute(
                    "DELETE FROM file_rounds WHERE session_id=? AND conversation_id=?",
                    (session_id, conversation_id)
                )
            elif session_id:
                # 删除整个会话的所有文件记录
                cursor.execute(
                    "DELETE FROM file_rounds WHERE session_id=?",
                    (session_id,)
                )
            
            self._db_conn.commit()
        except Exception as e:
            logger.error(f"删除文件使用次数记录失败: {str(e)}")
    
    def _get_session_id(self, event: AstrMessageEvent) -> str:
        """获取会话ID（统一消息来源）"""
        session_id = event.unified_msg_origin
        logger.debug(f"会话ID: {session_id}")
        return session_id
    
    async def _get_conversation_id(self, event: AstrMessageEvent) -> str:
        """获取当前对话ID"""
        session_id = self._get_session_id(event)
        conversation_id = await self.context.conversation_manager.get_curr_conversation_id(session_id)
        
        if not conversation_id:
            conversation_id = await self.context.conversation_manager.new_conversation(session_id)
            logger.info(f"为会话 {session_id} 创建新对话: {conversation_id}")
        else:
            logger.debug(f"使用现有对话ID: {conversation_id}")
        
        return conversation_id
    
    def _generate_timestamped_filename(self, original_file_name: str) -> str:
        """生成带时间戳的文件名（文件名_时间戳）"""
        # 确保只使用文件名，不包含路径
        original_file_name = os.path.basename(original_file_name)
        timestamp = int(time.time())
        return f"{original_file_name}_{timestamp}"
    
    def _parse_timestamped_filename(self, timestamped_name: str) -> tuple:
        """从带时间戳的文件名中解析原始文件名和时间戳"""
        # 确保只处理文件名部分
        timestamped_name = os.path.basename(timestamped_name)
        
        # 找到最后一个下划线，确保后面是数字时间戳
        last_underscore = timestamped_name.rfind('_')
        if last_underscore != -1:
            timestamp_str = timestamped_name[last_underscore + 1:]
            if timestamp_str.isdigit():
                original_name = timestamped_name[:last_underscore]
                return original_name, int(timestamp_str)
        return timestamped_name, None
    
    def _is_file_expired(self, session_id: str, conversation_id: str, file_name: str) -> bool:
        """检查单个文件是否过期（时间和轮数）"""
        # 从文件名解析时间戳
        _, upload_time = self._parse_timestamped_filename(file_name)
        if upload_time is None:
            # 如果无法解析时间戳，认为文件已过期
            return True
        
        # 检查时间是否过期
        retention_time = self.file_retention_time * 60  # 转换为秒
        current_time = time.time()
        time_expired = (current_time - upload_time) > retention_time
        
        # 检查轮数是否过期
        max_rounds = self.file_max_rounds
        current_rounds = self._get_file_rounds(session_id, conversation_id, file_name)
        rounds_expired = current_rounds >= max_rounds
        
        return time_expired or rounds_expired

    async def _cleanup_unauthorized_group_files(self):
        """清理非启用群聊的文件数据库"""
        try:
            # 使用集合存储已经处理过的会话ID，避免重复清理
            processed_sessions = set()
            
            # 直接遍历数据目录下的所有会话ID目录
            if self._data_dir.exists():
                for session_dir in self._data_dir.iterdir():
                    if session_dir.is_dir():
                        session_id = session_dir.name
                        
                        # 如果该会话已经处理过，跳过
                        if session_id in processed_sessions:
                            continue
                        
                        # 尝试从会话ID中提取群聊ID
                        group_id = None
                        
                        # 检查是否为群聊会话（格式：适配器名称:GroupMessage:12345678）
                        if "GroupMessage" in session_id:
                            # 格式：适配器名称:GroupMessage:12345678
                            try:
                                parts = session_id.split(":")
                                if len(parts) >= 3:
                                    group_id = parts[2]
                            except Exception:
                                pass
                        
                        # 如果能够提取到群聊ID，认为是群聊会话
                        if group_id:
                            # 如果群聊文件处理被禁用，清理该群聊会话
                            if not self.enable_group_file_processing:
                                logger.info(f"群聊文件处理已禁用，清理群聊会话 {session_id} 的所有文件")
                                await self.cleanup_all_session_files(session_id)
                                processed_sessions.add(session_id)
                                continue
                            
                            # 如果配置了群聊白名单，检查群聊ID是否在白名单中
                            if self.enabled_groups:
                                if str(group_id) not in [str(g) for g in self.enabled_groups]:
                                    logger.info(f"群聊 {group_id} 不在白名单中，清理会话 {session_id} 的所有文件")
                                    await self.cleanup_all_session_files(session_id)
                                    processed_sessions.add(session_id)
        except Exception as e:
            logger.error(f"清理非启用群聊文件失败: {str(e)}")
    
    async def initialize(self):
        """初始化嵌入提供者和重排序提供者"""
        try:
            # 首先重置提供者状态，确保每次初始化都是全新尝试
            self.embedding_provider = None
            self.rerank_provider = None
            
            # 使用类属性获取嵌入提供者ID
            embedding_provider_id = self.embedding_provider_id
            
            # 如果配置了特定的嵌入提供者ID，使用该提供者
            if embedding_provider_id:
                self.embedding_provider = self.context.get_provider_by_id(embedding_provider_id)
                logger.info(f"使用配置的嵌入提供者: {embedding_provider_id}")
            
            # 如果没有配置或者获取失败，使用默认的嵌入提供者
            if not self.embedding_provider:
                embedding_providers = self.context.get_all_embedding_providers()
                for provider in embedding_providers:
                    if hasattr(provider, 'get_embedding') and not self.embedding_provider:
                        self.embedding_provider = provider
                        break
            
            # 使用类属性获取重排序提供者ID
            rerank_provider_id = self.rerank_provider_id
            
            # 如果配置了特定的重排序提供者ID，使用该提供者
            if rerank_provider_id:
                self.rerank_provider = self.context.get_provider_by_id(rerank_provider_id)
                logger.info(f"使用配置的重排序提供者: {rerank_provider_id}")
            
            # 如果没有配置或者获取失败，使用默认的重排序提供者
            if not self.rerank_provider:
                # 直接从provider_manager获取所有重排序提供者
                rerank_providers = self.context.provider_manager.rerank_provider_insts
                for provider in rerank_providers:
                    if hasattr(provider, 'rerank') and not self.rerank_provider:
                        self.rerank_provider = provider
                        break
                
                # 如果直接访问provider_manager失败，尝试从所有提供者中过滤
                if not self.rerank_provider:
                    all_providers = self.context.provider_manager.inst_map.values()
                    for provider in all_providers:
                        if isinstance(provider, RerankProvider) and hasattr(provider, 'rerank') and not self.rerank_provider:
                            self.rerank_provider = provider
                            break
            
            if not self.embedding_provider:
                logger.error("无法获取嵌入提供者")
                return False
            
            logger.info(f"使用的嵌入提供者: {self.embedding_provider.__class__.__name__}")
            if self.rerank_provider:
                logger.info(f"使用的重排序提供者: {self.rerank_provider.__class__.__name__}")
            else:
                logger.warning("无法获取重排序提供者，将不使用重排序功能")
            
            # 清理非启用群聊的文件数据库
            await self._cleanup_unauthorized_group_files()
            
            # 启动定期清理任务
            await self._start_periodic_cleanup()
            return True
        except Exception as e:
            logger.error(f"初始化提供者失败: {str(e)}")
            return False
    
    async def get_or_create_vector_db(self, session_id: str, conversation_id: str, file_name: str):
        """获取或创建向量数据库（按会话、对话和文件名隔离）"""
        if not self.embedding_provider:
            logger.error("嵌入提供者未初始化，无法创建向量数据库")
            return None
        
        db_key = (session_id, conversation_id, file_name)
        
        # 如果已存在该会话/对话/文件的向量数据库，直接返回
        if db_key in self.vec_dbs:
            return self.vec_dbs[db_key]
        
        try:
            # 创建向量数据库目录（使用标准数据目录）
            vec_db_dir = self._data_dir / session_id / conversation_id / file_name
            vec_db_dir.mkdir(parents=True, exist_ok=True)
            
            # 初始化向量数据库
            vec_db = FaissVecDB(
                doc_store_path=str(vec_db_dir / "doc.db"),
                index_store_path=str(vec_db_dir / "index.faiss"),
                embedding_provider=self.embedding_provider,
                rerank_provider=self.rerank_provider
            )
            await vec_db.initialize()
            
            # 将向量数据库实例添加到字典中
            self.vec_dbs[db_key] = vec_db
            logger.info(f"为会话 {session_id} 对话 {conversation_id} 文件 {file_name} 创建向量数据库成功")
            return vec_db
        except Exception as e:
            logger.error(f"初始化向量数据库失败: {str(e)}")
            return None

    async def cleanup(self, session_id: str = None, conversation_id: str = None, file_name: str = None):
        """清理资源
        
        - 如果提供了session_id、conversation_id和file_name：清理单个文件
        - 如果提供了session_id和conversation_id：清理整个对话
        - 如果只提供了session_id：清理整个会话
        """
        if not session_id:
            session_id = self.current_session_id
        if not conversation_id:
            conversation_id = self.current_conversation_id
            
        if not session_id:
            logger.warning("未指定会话ID，无法清理")
            return
            
        if file_name:
            # 清理单个文件
            key = (session_id, conversation_id, file_name)
            if key in self.vec_dbs:
                vec_db = self.vec_dbs[key]
                await vec_db.close()
                
                # 删除向量数据库文件
                try:
                    vec_db_dir = self._data_dir / session_id / conversation_id / file_name
                    if vec_db_dir.exists():
                        import shutil
                        shutil.rmtree(vec_db_dir)
                except Exception as e:
                    logger.error(f"清理向量数据库文件失败: {str(e)}")
                
                # 从字典中移除
                del self.vec_dbs[key]
                
                # 从数据库中删除该文件的使用次数记录
                self._delete_file_rounds(session_id, conversation_id, file_name)
                
                logger.info(f"已清理会话 {session_id} 对话 {conversation_id} 的文件 {file_name} 向量数据库")
        elif conversation_id:
            # 清理整个对话
            keys_to_remove = []
            for (db_session_id, db_conversation_id, db_file_name), vec_db in self.vec_dbs.items():
                if db_session_id == session_id and db_conversation_id == conversation_id:
                    await vec_db.close()
                    keys_to_remove.append((db_session_id, db_conversation_id, db_file_name))
            
            # 从字典中移除
            for key in keys_to_remove:
                del self.vec_dbs[key]
            
            # 删除对话目录
            try:
                conversation_dir = self._data_dir / session_id / conversation_id
                if conversation_dir.exists():
                    import shutil
                    shutil.rmtree(conversation_dir)
            except Exception as e:
                logger.error(f"清理对话目录失败: {str(e)}")
            
            # 从数据库中删除该对话的所有文件使用次数记录
            self._delete_file_rounds(session_id, conversation_id)
            
            logger.info(f"已清理会话 {session_id} 对话 {conversation_id} 的所有文件")
        else:
            # 清理整个会话
            await self.cleanup_all_session_files(session_id)
            
            # 从数据库中删除该会话的所有文件使用次数记录
            self._delete_file_rounds(session_id)
        
    async def cleanup_all_session_files(self, session_id):
        """清理指定会话的所有文件"""
        try:
            # 关闭并删除该会话下的所有向量数据库实例
            keys_to_remove = []
            for (db_session_id, db_conversation_id, db_file_name), vec_db in self.vec_dbs.items():
                if db_session_id == session_id:
                    await vec_db.close()
                    keys_to_remove.append((db_session_id, db_conversation_id, db_file_name))
            
            # 从字典中移除已关闭的向量数据库实例
            for key in keys_to_remove:
                del self.vec_dbs[key]
            
            # 删除该会话下的所有对话文件
            session_dir = self._data_dir / session_id
            if session_dir.exists():
                import shutil
                shutil.rmtree(session_dir)
                logger.info(f"已清理会话 {session_id} 的所有文件")
        except Exception as e:
            logger.error(f"清理会话 {session_id} 的所有文件失败: {str(e)}")

    @filter.command("clear_file")
    async def clear_file_command(self, event: AstrMessageEvent):
        '''清理当前用户的所有文件''' 
        current_session_id = self._get_session_id(event)
        await self.cleanup_all_session_files(current_session_id)
        self.content = ""
        self.file_name = ""
        self.file_upload_time = None
        yield event.plain_result(f"已清理当前用户的所有文件，可以上传新文件了😊")

    @filter.command("clean_file")
    async def clean_file_command(self, event: AstrMessageEvent):
        '''清理当前用户的所有文件''' 
        current_session_id = self._get_session_id(event)
        await self.cleanup_all_session_files(current_session_id)
        self.content = ""
        self.file_name = ""
        self.file_upload_time = None
        yield event.plain_result(f"已清理当前用户的所有文件，可以上传新文件了😊")

    @filter.event_message_type(filter.EventMessageType.ALL)               # type: ignore
    async def on_receive_msg(self, event: AstrMessageEvent):
        """当获取到有文件时"""
        # 检查是否有新文件上传
        has_file = False
        for item in event.message_obj.message:
            if isinstance(item, Comp.File):
                has_file = True
                break
        
        # 只有当有文件时才处理
        if has_file:
            # 检查是否为群聊消息
            is_group_message = event.get_message_type() == MessageType.GROUP_MESSAGE
            
            # 如果是群聊消息，检查是否启用群聊文件处理
            if is_group_message and not self.enable_group_file_processing:
                logger.info(f"群聊文件处理已禁用，忽略来自会话 {self.current_session_id} 的文件")
                return
                
            # 如果是群聊消息，检查是否在白名单中
            if is_group_message:
                # 获取群聊ID
                group_id = event.get_group_id()
                
                # 检查群聊白名单
                # 参考llm_poke插件，确保所有ID都转为字符串比较
                if self.enabled_groups and group_id and str(group_id) not in [str(g) for g in self.enabled_groups]:
                    logger.info(f"群聊 {group_id} 不在白名单中，忽略文件处理")
                    return
            
            # 获取会话ID和对话ID
            self.current_session_id = self._get_session_id(event)
            self.current_conversation_id = await self._get_conversation_id(event)
            self.current_file_rounds = 0  # 重置使用轮数
            
            for item in event.message_obj.message:
                if isinstance(item, Comp.File):# 判断有无File组件
                    try:
                        file_path = await item.get_file() # 获取文件
                        file_dir, raw_file_name = os.path.split(file_path)
                        # 确保file_name只包含文件名，不包含路径
                        file_name = os.path.basename(raw_file_name)
                        
                        # 检查文件大小
                        max_file_size_bytes = self.max_file_size * 1024 * 1024  # 转换为字节
                        file_size = os.path.getsize(file_path)
                        if file_size > max_file_size_bytes:
                            logger.warning(f"文件 {file_name} 大小超过限制 ({file_size / 1024 / 1024:.2f}MB > {self.max_file_size}MB)")
                            yield event.plain_result(f"文件 {file_name} 大小超过限制 ({file_size / 1024 / 1024:.2f}MB > {self.max_file_size}MB)")
                            return
                        
                        # 获取完整文件名以确定正确的文件类型
                        completed_name = complete_filename(file_path)
                        # 检查文件类型是否支持
                        file_ext = os.path.splitext(completed_name)[1][1:].lower() if os.path.splitext(completed_name)[1] else ""
                        if file_ext and file_ext not in self.supported_file_types:
                            logger.warning(f"不支持的文件类型: {file_ext}")
                            yield event.plain_result(f"不支持的文件类型: {file_ext}")
                            return
                        
                        logger.info(f"接收到文件: {file_name}, 文件路径：{file_path}, 大小：{file_size / 1024 / 1024:.2f}MB")
                        # yield event.plain_result(f"已接收文件：{file_name}，正在处理...")
                        
                        # 读取文件内容
                        content = read_any_file_to_text(file_path)
                        
                        # 检查是否为错误信息
                        error_prefixes = ["文件不存在:", "不支持 ", "找不到处理 ", "读取文件时出错:"]
                        is_error = any(content.startswith(prefix) for prefix in error_prefixes)
                        
                        if content and not is_error:
                            logger.info(f"读取文件{file_name}内容成功")
                            
                            # 检查模型是否可用，如果不可用尝试重新获取
                            model_available = False
                            max_retries = 2
                            retry_count = 0
                            
                            while retry_count < max_retries:
                                # 检查嵌入提供者是否可用
                                if self.embedding_provider:
                                    logger.info("嵌入提供者已初始化，跳过重新获取")
                                    model_available = True
                                    break
                                else:
                                    logger.warning(f"嵌入提供者不可用，尝试重新获取 (第{retry_count + 1}次)")
                                    # 尝试重新初始化提供者
                                    init_success = await self.initialize()
                                    if init_success:
                                        logger.info("重新获取模型成功")
                                        model_available = True
                                        break
                                    else:
                                        retry_count += 1
                                        logger.error(f"重新获取模型失败，剩余重试次数: {max_retries - retry_count}")
                            
                            if model_available:
                                # 生成带时间戳的数据库名称
                                timestamped_db_name = self._generate_timestamped_filename(file_name)
                                
                                # 获取或创建向量数据库（需要会话、对话ID和带时间戳的文件名）
                                vec_db = await self.get_or_create_vector_db(self.current_session_id, self.current_conversation_id, timestamped_db_name)
                                
                                if vec_db:
                                    # 将文件内容分块
                                    chunks = await self.chunker.chunk(content)
                                    logger.info(f"文件分块完成，共{len(chunks)}个块")
                                    
                                    # 将块存入向量数据库
                                    metadatas = [{"file_name": file_name, "chunk_index": i} for i, _ in enumerate(chunks)]
                                    await vec_db.insert_batch(chunks, metadatas)
                                    logger.info(f"文件内容已存入向量数据库")
                                    logger.info(f"使用带时间戳的数据库名称：{timestamped_db_name}")

                                    # 成功向量化后，删除原始文件
                                    try:
                                        os.remove(file_path)
                                        logger.info(f"文件 {file_name} 已成功向量化并删除原始文件")
                                    except Exception as e:
                                        logger.warning(f"删除原始文件 {file_name} 失败: {str(e)}")

                                    yield event.plain_result(f"文件：{file_name} 已预处理完毕！请随时提问~ 😊")
                            else:
                                logger.error(f"无法获取可用的嵌入提供者，无法处理文件 {file_name}")
                                yield event.plain_result(f"文件处理失败：无法获取模型服务，请稍后重试或检查配置")
                        elif is_error:
                            logger.warning(f"读取文件{file_name}失败: {content}")
                            yield event.plain_result(content)  # 返回错误信息给用户
                        else:
                            logger.warning(f"读取文件{file_name}内容为空")
                    except Exception as e:
                        logger.error(f"读取文件失败: {str(e)}")

    @filter.on_llm_request(proirity=-9999)
    async def on_request(self, event: AstrMessageEvent, req: ProviderRequest):
        # 获取当前会话和对话ID
        current_session_id = self._get_session_id(event)
        current_conversation_id = await self._get_conversation_id(event)
        
        # 更新当前会话和对话ID
        if (current_session_id != self.current_session_id or 
            current_conversation_id != self.current_conversation_id):
            self.current_session_id = current_session_id
            self.current_conversation_id = current_conversation_id
        
        # 获取当前会话/对话下的所有文件向量数据库
        all_results_with_source = []
        all_files = set()
        
        # 遍历所有向量数据库，检查是否属于当前会话/对话
        for (db_session_id, db_conversation_id, file_name), vec_db in list(self.vec_dbs.items()):
            if db_session_id == current_session_id and db_conversation_id == current_conversation_id:
                # 检查文件是否过期
                if self._is_file_expired(db_session_id, db_conversation_id, file_name):
                    logger.info(f"文件 {file_name} 已过期，将清理并停止使用")
                    await self.cleanup(db_session_id, db_conversation_id, file_name)
                    continue
                
                # 解析出原始文件名用于显示（从实际访问的数据库路径获取）
                original_file_name, _ = self._parse_timestamped_filename(file_name)
                all_files.add(original_file_name)
                logger.info(f"从文件 {original_file_name} 的向量数据库检索与查询相关的内容")
                
                # 从请求中获取用户查询
                user_query = req.prompt
                
                # 使用类属性进行检索
                retrieve_top_k = self.retrieve_top_k
                fetch_k = self.fetch_k
                enable_rerank = self.enable_rerank
                
                # 检索相关内容
                results = await vec_db.retrieve(user_query, k=retrieve_top_k, fetch_k=fetch_k, rerank=enable_rerank)
                
                # 记录每个结果来自哪个数据库文件
                for result in results:
                    all_results_with_source.append((result, original_file_name))
        
        if all_results_with_source:
            logger.info(f"共检索到{len(all_results_with_source)}条相关内容")
            
            # 构建上下文
            context_text = "以下是与查询相关的文件内容:\n"
            
            # 添加相关文件列表
            if all_files:
                context_text += f"相关文件: {', '.join(all_files)}\n\n"
            
            # 添加相关内容
            for i, (result, file_name) in enumerate(all_results_with_source, 1):
                # 确保result.data是字典
                if hasattr(result, 'data') and isinstance(result.data, dict):
                    chunk = result.data.get("text", "")
                    context_text += f"\n【文件: {file_name} 片段{i}】\n{chunk}\n"
            
            # 根据配置选择注入方式
            if self.injection_type == "system":
                # 清理系统上下文，保留最后N-1轮的system字段
                # 因为当前请求也会算作一轮，所以实际保留的是N-1轮
                if self.system_context_keep_rounds > 0:
                    # 如果keep_rounds=1，清理所有system消息
                    if self.system_context_keep_rounds == 1:
                        # 清理所有system消息
                        new_contexts = [ctx for ctx in contexts if ctx.get("role") != "system"]
                        req.contexts = new_contexts
                        logger.debug(f"已清理所有系统上下文，保留最近 {self.system_context_keep_rounds-1} 轮的完整内容")
                    else:
                        # 简化的清理逻辑：找到倒数第keep_rounds轮的结束位置，清除之前的所有system消息
                        # keep_rounds=2时找倒数第2轮的结束位置，keep_rounds=3时找倒数第3轮的结束位置，以此类推
                        contexts = req.contexts
                        new_contexts = []
                        
                        # 从后往前找，找到倒数第keep_rounds轮的结束位置（assistant）
                        cutoff_index = -1  # 默认不清除任何system消息
                        i = len(contexts) - 1
                        rounds_found = 0
                        
                        while i >= 0 and rounds_found < self.system_context_keep_rounds:
                            # 找到当前轮次的最后一个assistant
                            while i >= 0 and contexts[i].get("role") != "assistant":
                                i -= 1
                            
                            if i < 0:
                                break
                                
                            # 找到了一个assistant，这是一个轮次的结束
                            rounds_found += 1
                            
                            # 如果这是我们需要找的轮次（倒数第keep_rounds轮），记录其位置
                            if rounds_found == self.system_context_keep_rounds:
                                cutoff_index = i
                                break
                            
                            # 继续往前找下一个轮次
                            # 跳过当前轮次的所有内容，直到找到上一个轮次的user或system
                            while i >= 0 and contexts[i].get("role") not in ["user", "system"]:
                                i -= 1
                            
                            # 跳过这个user或system
                            if i >= 0:
                                i -= 1
                        
                        # 遍历所有上下文，决定是否保留
                        for i, ctx in enumerate(contexts):
                            role = ctx.get("role")
                            
                            # 如果是system消息且在cutoff_index之前，则清除
                            if role == "system" and i <= cutoff_index:
                                continue  # 跳过这个system消息
                            else:
                                new_contexts.append(ctx)
                        
                        # 更新上下文
                        req.contexts = new_contexts
                        logger.debug(f"已清理系统上下文，保留最近 {self.system_context_keep_rounds-1} 轮的完整内容")
                
                # 将文件内容注入到系统上下文
                system_prompt = f"文件相关内容:\n{context_text}\n\n请根据上述内容回答用户问题:"
                req.contexts.append({"role": "system", "content": system_prompt})
                # 保持原始用户查询作为prompt
                req.prompt = user_query
                logger.info(f"已将文件内容以system类型注入到请求中")
            else:
                # 将文件内容注入到用户prompt中（默认行为）
                req.prompt = f"{user_query}\n\n文件相关内容:\n{context_text}\n\n请根据上述内容回答用户问题:"
                logger.info(f"已将文件内容以user类型注入到请求中")
        elif all_files:
            logger.info("未检索到相关内容")
        
        # 遍历当前会话/对话下的所有文件，为每个文件增加使用轮数
        for (db_session_id, db_conversation_id, db_file_name), _ in list(self.vec_dbs.items()):
            if db_session_id == current_session_id and db_conversation_id == current_conversation_id:
                self._increment_file_rounds(db_session_id, db_conversation_id, db_file_name)

    def __del__(self):
        """对象销毁时清理资源"""
        # 停止定期清理任务
        if hasattr(self, '_cleanup_task') and self._cleanup_task:
            self._cleanup_task.cancel()
            logger.info("已取消定期清理任务")
        
        # 清理资源 - 在__del__中避免使用异步操作，直接处理简单的资源释放
        # 更复杂的清理应该在对象正常使用时通过调用cleanup()方法完成
        for key, vec_db in list(self.vec_dbs.items()):
            try:
                # 尝试关闭向量数据库连接
                if hasattr(vec_db, 'close'):
                    vec_db.close()
            except Exception as e:
                logger.error(f"关闭向量数据库时出错: {str(e)}")
        self.vec_dbs.clear()